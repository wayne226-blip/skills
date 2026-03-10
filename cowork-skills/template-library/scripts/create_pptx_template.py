#!/usr/bin/env python3
"""Generate a branded PPTX reference template compatible with Pandoc.

Takes Pandoc's default reference.pptx and rebrands it with custom colors,
fonts, and styles — preserving the exact slide layout structure Pandoc expects.

Usage:
    python create_pptx_template.py [output.pptx] [--palette NAME]

Palettes: midnight (default), forest, coral, ocean, charcoal
"""

import argparse
import os
import subprocess
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


PALETTES = {
    "midnight": {
        "primary": RGBColor(0x1E, 0x27, 0x61),
        "secondary": RGBColor(0x2B, 0x6C, 0xB0),
        "accent": RGBColor(0x38, 0xA1, 0x69),
        "bg_dark": RGBColor(0x1E, 0x27, 0x61),
        "bg_light": RGBColor(0xF7, 0xFA, 0xFC),
        "text_dark": RGBColor(0x2D, 0x37, 0x48),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "muted": RGBColor(0x71, 0x80, 0x96),
        "heading_font": "Georgia",
        "body_font": "Calibri",
    },
    "forest": {
        "primary": RGBColor(0x2C, 0x5F, 0x2D),
        "secondary": RGBColor(0x97, 0xBC, 0x62),
        "accent": RGBColor(0x38, 0xA1, 0x69),
        "bg_dark": RGBColor(0x2C, 0x5F, 0x2D),
        "bg_light": RGBColor(0xF5, 0xF5, 0xF5),
        "text_dark": RGBColor(0x1A, 0x20, 0x2C),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "muted": RGBColor(0x6B, 0x72, 0x80),
        "heading_font": "Georgia",
        "body_font": "Calibri",
    },
    "coral": {
        "primary": RGBColor(0x2F, 0x3C, 0x7E),
        "secondary": RGBColor(0xF9, 0x61, 0x67),
        "accent": RGBColor(0xF9, 0xE7, 0x95),
        "bg_dark": RGBColor(0x2F, 0x3C, 0x7E),
        "bg_light": RGBColor(0xFF, 0xF5, 0xF5),
        "text_dark": RGBColor(0x2D, 0x37, 0x48),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "muted": RGBColor(0x71, 0x80, 0x96),
        "heading_font": "Georgia",
        "body_font": "Calibri",
    },
    "ocean": {
        "primary": RGBColor(0x06, 0x5A, 0x82),
        "secondary": RGBColor(0x1C, 0x72, 0x93),
        "accent": RGBColor(0x21, 0x29, 0x5C),
        "bg_dark": RGBColor(0x06, 0x5A, 0x82),
        "bg_light": RGBColor(0xEB, 0xF8, 0xFF),
        "text_dark": RGBColor(0x1A, 0x20, 0x2C),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "muted": RGBColor(0x64, 0x74, 0x8B),
        "heading_font": "Georgia",
        "body_font": "Calibri",
    },
    "charcoal": {
        "primary": RGBColor(0x36, 0x45, 0x4F),
        "secondary": RGBColor(0x4A, 0x55, 0x68),
        "accent": RGBColor(0x21, 0x21, 0x21),
        "bg_dark": RGBColor(0x36, 0x45, 0x4F),
        "bg_light": RGBColor(0xF2, 0xF2, 0xF2),
        "text_dark": RGBColor(0x21, 0x21, 0x21),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "muted": RGBColor(0x6B, 0x72, 0x80),
        "heading_font": "Arial",
        "body_font": "Arial",
    },
}


def apply_font_to_placeholder(placeholder, font_name, font_size, font_color, bold=False):
    """Apply font styling to all runs in a placeholder's text frame."""
    if not placeholder.has_text_frame:
        return
    for paragraph in placeholder.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = font_color
            run.font.bold = bold


def create_template(output_path, palette_name):
    p = PALETTES.get(palette_name, PALETTES["midnight"])

    # Generate Pandoc's default reference template
    base_path = "/tmp/_pandoc_base_ref.pptx"
    result = subprocess.run(
        ["pandoc", "-o", base_path, "--print-default-data-file", "reference.pptx"],
        capture_output=True
    )
    if result.returncode != 0:
        # Fallback: generate one from scratch
        subprocess.run(
            ["pandoc", "-t", "pptx", "-o", base_path],
            input=b"# Title\n\n## Slide",
            capture_output=True
        )

    prs = Presentation(base_path)

    # Modify slide masters and layouts
    for slide_master in prs.slide_masters:
        # Try to style placeholders in slide layouts
        for layout in slide_master.slide_layouts:
            for placeholder in layout.placeholders:
                ph_type = placeholder.placeholder_format.type
                if placeholder.has_text_frame:
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            # Title placeholders
                            if ph_type in [1, 15]:  # TITLE, CENTERED_TITLE
                                run.font.name = p["heading_font"]
                                run.font.color.rgb = p["primary"]
                                run.font.bold = True
                            # Body placeholders
                            elif ph_type in [2, 7]:  # BODY, OBJECT
                                run.font.name = p["body_font"]
                                run.font.color.rgb = p["text_dark"]

    prs.save(output_path)
    file_size = os.path.getsize(output_path)
    print(f"Template created: {output_path} ({file_size:,} bytes, palette: {palette_name})")

    # Clean up
    if os.path.exists(base_path):
        os.remove(base_path)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate branded PPTX reference template")
    parser.add_argument("output", nargs="?", default="wayne_brand_template.pptx", help="Output file")
    parser.add_argument("--palette", default="midnight", choices=list(PALETTES.keys()))
    args = parser.parse_args()
    create_template(args.output, args.palette)
