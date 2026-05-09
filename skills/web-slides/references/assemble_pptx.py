"""Assemble slide PNGs into a PPTX file using python-pptx.

Usage:
    python3 assemble_pptx.py <png-dir/> <output.pptx>

Each PNG is added as a full-bleed image slide at 16:9 (1920x1080).
PNGs are sorted by filename (slide-01.png, slide-02.png, etc.).
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    print("ERROR: python-pptx not installed.", file=sys.stderr)
    print("Run: pip3 install python-pptx", file=sys.stderr)
    sys.exit(1)


def assemble(png_dir: Path, output_path: Path) -> Path:
    """Create a PPTX with each PNG as a full-bleed slide. Returns output path."""
    pngs = sorted(png_dir.glob("slide-*.png"))
    if not pngs:
        print(f"ERROR: No slide-*.png files found in {png_dir}", file=sys.stderr)
        sys.exit(1)

    # 16:9 slide dimensions: 13.333" x 7.5" (standard widescreen)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Use a blank layout
    blank_layout = prs.slide_layouts[6]  # Blank slide layout

    for png_path in pngs:
        slide = prs.slides.add_slide(blank_layout)
        # Add image at full-bleed (0,0 position, full slide dimensions)
        slide.shapes.add_picture(
            str(png_path),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        print(f"  Added {png_path.name}")

    prs.save(str(output_path))
    print(f"\nDone. {len(pngs)} slides -> {output_path}")
    return output_path


def main() -> None:
    parser = argparse.ArgumentParser(description="Assemble slide PNGs into PPTX")
    parser.add_argument("png_dir", type=Path, help="Directory containing slide-*.png files")
    parser.add_argument("output", type=Path, help="Output PPTX file path")
    args = parser.parse_args()

    if not args.png_dir.is_dir():
        print(f"ERROR: Directory not found: {args.png_dir}", file=sys.stderr)
        sys.exit(1)

    assemble(args.png_dir, args.output)


if __name__ == "__main__":
    main()
